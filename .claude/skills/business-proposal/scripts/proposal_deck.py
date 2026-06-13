# -*- coding: utf-8 -*-
"""
綠芽 GRASSIDEA 企業形象影片提案企劃 — 簡報產生器（對標鼎康範例視覺系統）

定案格式：每個切入點 3 頁（區塊頁 / 核心訊息主視覺 / 影片架構＋內容大綱）。
全程 python-pptx 原生繪製，照片區一律用色塊（禁用 JPG/PNG）。

用法：
    from proposal_deck import build_deck
    build_deck(
        client_name="品香茶葉",
        angles=[ {...}, {...}, {...} ],   # 見下方 ANGLE schema
        out_path="/path/【品香茶葉】企業形象提案企劃.pptx",
    )

ANGLE schema（每個切入點一個 dict）：
    no:        "01" / "02" / "03"
    name:      切入點名稱（如「可被驗證的信任」）
    hero_zh:   主視覺中文大標，用 \\n 斷行
    slogan:    英文 Slogan（broadcast 等級）
    sub:       英文副標
    statement: 一句中文核心陳述
    concepts:  3 個 (中文概念, 英文, 一句說明) tuple，對應「影片架構」概念圈
    labels:    架構底線標籤，如 ["Traceable", "Reliable"]
    outline:   內容大綱段落
    film:      影像方向一句話（導演 treatment）
    fit:       「適合：…」一句話
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 視覺系統（對標鼎康範例）----------
NAVY    = RGBColor(0x0C, 0x23, 0x40)   # 深藍：深色欄、標籤底
BLUE    = RGBColor(0x15, 0x65, 0xC0)   # 區塊頁主藍、accent
LIGHTBL = RGBColor(0xDB, 0xEA, 0xF7)   # 淺藍面板
GREEN   = RGBColor(0x2E, 0xD1, 0x5A)   # 綠芽品牌綠（僅封面/結尾）
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1A, 0x1A, 0x1A)   # 內文黑
GREY    = RGBColor(0x8A, 0x9B, 0xA8)   # 輔助灰
NEARBLK = RGBColor(0x0A, 0x0A, 0x0A)

EA, LAT = "PingFang TC", "Arial"        # 中文 / 英文字型
CN_NUM = "一二三四五六七八九十"


def _set_font(run, size=None, bold=None, color=None, italic=None, latin=LAT, ea=EA):
    f = run.font
    if size is not None:   f.size = Pt(size)
    if bold is not None:   f.bold = bold
    if italic is not None: f.italic = italic
    if color is not None:  f.color.rgb = color
    f.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag, face in (('a:latin', latin), ('a:ea', ea), ('a:cs', latin)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', face)


def _bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color


def _rect(s, x, y, w, h, color, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if color is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def _tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, color, bold=None, italic=None, align=PP_ALIGN.LEFT,
       latin=LAT, ea=EA, space_after=None, line=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None: p.space_after = Pt(space_after)
    if line is not None: p.line_spacing = line
    r = p.add_run(); r.text = text
    _set_font(r, size=size, bold=bold, color=color, italic=italic, latin=latin, ea=ea)
    return p


def _dot_grid(s, x0, y0, cols, rows, gap, size, color):
    for r in range(rows):
        for c in range(cols):
            _rect(s, x0 + c*gap, y0 + r*gap, size, size, color)


def _header(s, color=BLUE, text=" CORPORATE IMAGE VIDEO"):
    tf = _tb(s, Inches(8.5), Inches(0.28), Inches(4.6), Inches(0.35))
    _p(tf, text, 11, color, bold=True, align=PP_ALIGN.RIGHT, first=True)


def build_deck(client_name, angles, out_path, context=None,
               contact=("masenc13@grassidea13.com", "Masen Cheung", "+886 02 2506 2727")):
    """context（選填）= dict(summary=客戶一句話, audience=目標受眾, insights=[3 點市場洞察])
    有提供就在簡報大綱後插一頁「專案理解 · 市場洞察」，讓市場分析師的貢獻被看見。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    BLANK = prs.slide_layouts[6]
    EMU_W, EMU_H = prs.slide_width, prs.slide_height
    new = lambda: prs.slides.add_slide(BLANK)

    # 1. 封面
    s = new(); _bg(s, NEARBLK)
    tf = _tb(s, Inches(1.0), Inches(2.5), Inches(7), Inches(2.2))
    _p(tf, "GRASS", 54, WHITE, bold=True, first=True, line=1.0)
    _p(tf, "IDEA", 54, GREEN, bold=True, line=1.0)
    _p(tf, "STUDIO", 18, WHITE, bold=True)
    tf = _tb(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.2))
    _p(tf, f"{client_name}  企業形象影片提案企劃", 26, WHITE, bold=True, first=True)
    _p(tf, "Corporate Image Video — Proposal", 14, GREEN, italic=True)
    tf = _tb(s, Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.4))
    _p(tf, "          ".join(contact), 11, GREY, align=PP_ALIGN.CENTER, first=True)

    # 2. 簡報大綱
    s = new(); _bg(s, WHITE)
    _rect(s, 0, 0, Inches(0.55), EMU_H, NAVY)
    _rect(s, Inches(6.6), 0, Inches(6.73), EMU_H, LIGHTBL)
    _header(s)
    tf = _tb(s, Inches(2.2), Inches(2.0), Inches(4), Inches(2.4))
    _p(tf, "簡報", 60, NAVY, bold=True, first=True, line=1.0)
    _p(tf, "大綱", 60, NAVY, bold=True, line=1.0)
    y = Inches(2.0)
    for a in angles:
        tf = _tb(s, Inches(6.95), y, Inches(1.2), Inches(1.4))
        _p(tf, a["no"], 40, RGBColor(0xB8, 0xCF, 0xE6), bold=True, first=True)
        tf2 = _tb(s, Inches(8.05), y + Inches(0.18), Inches(5.0), Inches(1.3))
        _p(tf2, f"提案策略{CN_NUM[int(a['no'])-1]} · {a['name']}", 17, NAVY, bold=True, first=True)
        _p(tf2, a.get("slogan", ""), 12, GREY, italic=True)
        y += Inches(1.55)

    # 2.5 專案理解 · 市場洞察（市場分析師的貢獻，有 context 才出）
    if context:
        s = new(); _bg(s, WHITE)
        _rect(s, 0, 0, Inches(0.35), EMU_H, BLUE)
        _rect(s, Inches(12.98), 0, Inches(0.35), EMU_H, NAVY)
        _header(s)
        # —— 專案理解：淺藍資訊面板 ——
        lab = _rect(s, Inches(0.85), Inches(0.62), Inches(2.0), Inches(0.55), NAVY)
        _p(lab.text_frame, "專案理解", 17, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
        panel = _rect(s, Inches(0.85), Inches(1.42), Inches(11.63), Inches(1.5), LIGHTBL,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        try: panel.adjustments[0] = 0.06
        except Exception: pass
        ptf = panel.text_frame; ptf.word_wrap = True; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.margin_left = ptf.margin_right = Inches(0.45)
        _p(ptf, context.get("summary", ""), 16, INK, first=True, line=1.3, space_after=10)
        _p(ptf, "目標受眾　" + context.get("audience", ""), 14, BLUE, bold=True)
        # —— 市場洞察：3 張對齊卡片（數字圓 + 內文），做出圖表感 ——
        lab = _rect(s, Inches(0.85), Inches(3.32), Inches(2.0), Inches(0.55), NAVY)
        _p(lab.text_frame, "市場洞察", 17, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
        ins = context.get("insights", [])
        n = max(len(ins), 1); gap = Inches(0.4)
        cw = (Inches(11.63) - gap*(n-1)) / n
        cx, cy, chh = Inches(0.85), Inches(4.18), Inches(2.5)
        for i, txt in enumerate(ins):
            card = _rect(s, cx, cy, cw, chh, WHITE, line=RGBColor(0xC7, 0xD6, 0xE6),
                         line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            try: card.adjustments[0] = 0.05
            except Exception: pass
            d = Inches(0.74)                                    # 頂部藍色數字圓
            circ = _rect(s, cx + cw/2 - d/2, cy + Inches(0.3), d, d, BLUE, shape=MSO_SHAPE.OVAL)
            circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            _p(circ.text_frame, str(i+1), 24, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
            t = _tb(s, cx + Inches(0.28), cy + Inches(1.22), cw - Inches(0.56), Inches(1.1))
            _p(t, txt, 12.5, INK, align=PP_ALIGN.CENTER, first=True, line=1.26)
            cx += cw + gap

    # 每個切入點 3 頁
    for a in angles:
        # (a) 區塊頁
        s = new(); _bg(s, BLUE)
        _dot_grid(s, Inches(7.4), Inches(1.4), 11, 9, Inches(0.42), Inches(0.09), NAVY)
        _header(s, color=WHITE)
        tf = _tb(s, Inches(0.9), Inches(1.5), Inches(2), Inches(0.8))
        _p(tf, a["no"], 34, WHITE, bold=True, first=True)
        tf = _tb(s, Inches(0.85), Inches(2.8), Inches(6.2), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
        _p(tf, "提案策略" + CN_NUM[int(a["no"])-1], 30, WHITE, bold=True, first=True, line=1.05)
        _p(tf, a["name"], 44, WHITE, bold=True, line=1.05)
        tf = _tb(s, Inches(0.9), Inches(6.4), Inches(7), Inches(0.5))
        _p(tf, "核心訊息  ▪  影片架構  ▪  內容大綱", 15, WHITE, first=True)

        # (b) 核心訊息主視覺
        s = new(); _bg(s, NAVY)
        _rect(s, 0, Inches(4.7), EMU_W, Inches(2.8), RGBColor(0x07, 0x18, 0x2E))
        _header(s, color=WHITE)
        if a.get("image_hint"):   # 給後製/設計：這頁主視覺該放什麼畫面
            ht = _tb(s, Inches(0.9), Inches(0.72), Inches(7.4), Inches(0.4))
            _p(ht, "［ 主視覺示意：" + a["image_hint"] + " ］", 11,
               RGBColor(0x7E, 0x97, 0xB2), italic=True, first=True)
        tf = _tb(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(2.6))
        for i, ln in enumerate(a["hero_zh"].split("\n")):
            _p(tf, ln, 34, WHITE, bold=True, first=(i == 0), line=1.1)
        _p(tf, a["slogan"], 46, WHITE, italic=True, line=1.05)
        tf = _tb(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6))
        _p(tf, a["sub"], 17, RGBColor(0xC9, 0xD6, 0xE3), italic=True, first=True, line=1.25)
        _p(tf, a["statement"], 16, GREEN, bold=True)

        # (c) 影片架構 ＋ 內容大綱
        s = new(); _bg(s, WHITE)
        _rect(s, 0, 0, Inches(0.35), EMU_H, BLUE)
        _rect(s, Inches(12.98), 0, Inches(0.35), EMU_H, NAVY)
        _header(s)
        lab = _rect(s, Inches(5.55), Inches(0.55), Inches(2.2), Inches(0.62), NAVY)
        _p(lab.text_frame, "影片架構", 18, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
        n = len(a["concepts"]); gap = Inches(0.4)
        cw = (Inches(11.6) - gap*(n-1)) / n
        x, yy, ch = Inches(0.85), Inches(1.55), Inches(2.0)
        for (zh, en, desc) in a["concepts"]:
            pill = _rect(s, x, yy, cw, ch, WHITE, line=NAVY, line_w=Pt(1.25),
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            try: pill.adjustments[0] = 0.18
            except Exception: pass
            ptf = pill.text_frame; ptf.word_wrap = True
            ptf.vertical_anchor = MSO_ANCHOR.TOP
            ptf.margin_left = ptf.margin_right = Inches(0.25); ptf.margin_top = Inches(0.30)
            _p(ptf, zh, 19, NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
            _p(ptf, en, 11, BLUE, bold=True, align=PP_ALIGN.CENTER, space_after=6)
            _p(ptf, desc, 12, INK, align=PP_ALIGN.CENTER, line=1.2)
            x += cw + gap
        ly = Inches(3.82)
        ln = s.shapes.add_connector(2, Inches(0.85), ly, Inches(12.45), ly)
        ln.line.color.rgb = GREY; ln.line.width = Pt(0.75)
        tf = _tb(s, Inches(0.85), ly + Inches(0.06), Inches(11.6), Inches(0.4))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "      ".join(a["labels"])
        _set_font(r, size=18, color=BLUE)
        # 提案重點：市場依據（市場分析師）→ 內容大綱 → 影像方向（導演）→ 適合
        olab = _rect(s, Inches(0.85), Inches(4.45), Inches(0.62), Inches(2.2), NAVY)
        otf = olab.text_frame; otf.word_wrap = True; otf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _p(otf, "提\n案\n重\n點", 15, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, line=1.0)
        tf = _tb(s, Inches(1.75), Inches(4.45), Inches(10.6), Inches(2.3))
        _p(tf, "市場依據　" + a.get("market", ""), 12.5, BLUE, bold=True, first=True,
           line=1.24, space_after=6)
        _p(tf, a["outline"], 13, INK, line=1.26, space_after=6)
        _p(tf, a["film"], 12, RGBColor(0x33, 0x55, 0x77), line=1.22, space_after=4)
        _p(tf, a["fit"], 12, NAVY, bold=True, line=1.22)

    # 結尾 THANK YOU
    s = new(); _bg(s, NEARBLK)
    tf = _tb(s, Inches(0.9), Inches(0.7), Inches(4), Inches(1.2))
    _p(tf, "GRASS IDEA", 20, WHITE, bold=True, first=True)
    _p(tf, "STUDIO", 11, GREEN, bold=True)
    tf = _tb(s, Inches(2.0), Inches(2.7), Inches(9.3), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
    _p(tf, "“ THANK YOU. ”", 66, GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)
    tf = _tb(s, Inches(0.9), Inches(5.1), Inches(6), Inches(0.9))
    _p(tf, contact[1], 15, WHITE, bold=True, first=True)
    _p(tf, contact[0], 13, GREY)
    tf = _tb(s, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4))
    _p(tf, "www.grassidea13.com", 11, GREEN, first=True)

    prs.save(out_path)
    return len(prs.slides._sldIdLst)


if __name__ == "__main__":
    # 範例：品香茶葉（定案內容，可作為新案的填寫範本）
    angles = [
        dict(no="01", name="可被驗證的信任",
             hero_zh="看得見的安心，\n賣得出去的信任",
             slogan="Every Leaf, Fully Accountable.",
             sub="From Our Soil to Your Container — Traceable at Every Step.",
             statement="自耕自產，從一片茶葉到一只貨櫃，全程由我們把關。",
             image_hint="黑底茶菁 slow-mo 落下、製程微距特寫、職人之手",
             market="海外買家最大風險是農殘超標退關；自耕自產＝單一責任主體，是最硬的進場門票。",
             concepts=[("源頭掌控", "Single Origin", "自有茶園，從茶籽到茶菁全程自管，杜絕來源不明的中間商。"),
                       ("工藝可見", "Craft on Display", "製茶工序與 ESG 封裝以 demonstration 鏡頭呈現專業與食安。"),
                       ("安心出口", "Ready to Ship", "代工封裝、上棧板待運，穩定供貨、責任可問。")],
             labels=["Traceable", "Reliable"],
             outline="以「可被驗證」為主軸，用製程與封裝畫面當證據，讓海外採購一眼看懂品香的食安與供貨保證。",
             film="影像方向：深色控光、probe lens 微距、slow motion 製程特寫，聚焦職人之手。",
             fit="適合：主攻歐美日嚴檢市場、需要快速建立進場信任時。"),
        dict(no="02", name="台灣風土，不可複製",
             hero_zh="一片茶葉，\n一座島嶼的風土",
             slogan="Tea That Could Only Be Taiwan.",
             sub="Mountains, Mist, and the Hands That Know Them.",
             statement="我們賣的不是茶，是只有台灣土地長得出來的味道。",
             image_hint="晨霧茶山 drone slow push-in、採茶人剪影、露水茶菁特寫",
             market="對標葡萄酒 terroir，買家可轉售故事創造溢價；差異化對手是日本。",
             concepts=[("風土之源", "Terroir", "高山、雲霧、氣候，台灣獨有、不可複製的風味地景。"),
                       ("人文之手", "Heritage", "世代製茶職人，風土與工藝的傳承共同養出這杯茶。"),
                       ("風味之證", "Signature", "MADE IN TAIWAN 從產地標示，升級為品質宣言。")],
             labels=["Emotional", "Distinctive"],
             outline="以風土詩意敘事，把台灣高山的自然與人文轉成不可複製的品牌資產，讓 MADE IN TAIWAN 成為情感高點。",
             film="影像方向：自然紀實、霧綠大地色調、magic hour 空拍緩推，舒緩呼吸式剪輯。",
             fit="適合：要建立長線品牌資產、走精品溢價路線時。"),
        dict(no="03", name="一站式外銷夥伴",
             hero_zh="從一座島，\n到世界的杯中",
             slogan="From One Island, to the World's Cup.",
             sub="We Grow It, Craft It, Pack It, and Ship It.",
             statement="能種、能製、能代工、能出貨——你的長期茶葉供應夥伴。",
             image_hint="地球曲線晨光俯衝穿雲落到台灣茶山、棧板待運 wide",
             market="買家要供貨彈性與 OEM/ODM 能力，要的是能長期配合的供應商。",
             concepts=[("一條龍實力", "End-to-End", "自耕、製茶、封裝、出口一站整合，供貨穩定。"),
                       ("客製彈性", "OEM / ODM", "配合代工與客製包裝需求，彈性對應通路規格。"),
                       ("全球視野", "Global Ready", "以外銷企圖與品牌格局，走向世界市場。")],
             labels=["Flexible", "Scalable"],
             outline="以品牌大器的國際視野，展現品香從風土到全球供應鏈的一條龍實力，爭取長期代工與通路合作。",
             film="影像方向：cinematic brand film、anamorphic 寬幅、teal-and-warm 高級調，大器運鏡。",
             fit="適合：要展現外銷企圖與供應鏈實力、爭取長期合作時。"),
    ]
    context = dict(
        summary="品香茶葉：台灣自耕自產的茶葉品牌，核心訴求「優質好茶」，主打天然、自耕、來自台灣、食品安全。",
        audience="海外茶葉進口商／通路採購（外銷導向）",
        insights=["買家第一痛點是農殘超標退關——產地信任＝風險評估，影片必須先建立安心感。",
                  "台灣茶走「小而精」利基，差異化對手是中國（拚量）與日本（拚形象），打的是可信任的精品產地。",
                  "自耕自產＝源頭到成品單一責任主體，是品香最強、競品難複製的 B2B 信任資產。"])
    out = "/Users/masen/Desktop/執行中案件/品香茶葉_wendy/【品香茶葉】企業形象提案企劃.pptx"
    print("頁數：", build_deck("品香茶葉", angles, out, context=context))
    print("已儲存：", out)
