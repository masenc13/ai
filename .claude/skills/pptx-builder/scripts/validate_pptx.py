#!/usr/bin/env python3
"""
PPTX 版面驗證工具
用法：python3 validate_pptx.py <path_to_pptx>
"""
import sys
from pptx import Presentation
from pptx.util import Inches

TOLERANCE = 0.02  # 允許 0.02" 誤差

def is_watermark_number(text: str) -> bool:
    """刻意放的深色浮水印章節數字"""
    return text.strip() in {"01", "02", "03", "04", "05", "06", "07", "08", "09",
                             "1", "2", "3", "4", "5", "6", "7", "8", "9", "10"}

def is_footer(text: str) -> bool:
    return text.strip().startswith("©")

def get_textboxes(slide):
    result = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text.strip()
        if not text:
            continue
        l  = sh.left.inches   if sh.left    else 0
        t  = sh.top.inches    if sh.top     else 0
        w  = sh.width.inches  if sh.width   else 0
        h  = sh.height.inches if sh.height  else 0
        result.append({
            "name": sh.name,
            "text": text,
            "l": l, "t": t, "w": w, "h": h,
            "r": l + w, "b": t + h,
        })
    return result

def check_slide(slide, slide_w, slide_h, page_num):
    issues = []
    tbs = get_textboxes(slide)

    # 1. Overflow
    for tb in tbs:
        if tb["r"] > slide_w + TOLERANCE:
            issues.append({
                "type": "overflow_x",
                "page": page_num,
                "name": tb["name"],
                "text": tb["text"][:40],
                "detail": f"right={tb['r']:.3f}\" > slide_width={slide_w:.3f}\"",
            })
        if tb["b"] > slide_h + TOLERANCE:
            issues.append({
                "type": "overflow_y",
                "page": page_num,
                "name": tb["name"],
                "text": tb["text"][:40],
                "detail": f"bottom={tb['b']:.3f}\" > slide_height={slide_h:.3f}\"",
            })

    # 2. 文字重疊
    for i in range(len(tbs)):
        for j in range(i + 1, len(tbs)):
            a, b = tbs[i], tbs[j]
            # 排除浮水印數字 & footer
            if is_watermark_number(a["text"]) or is_watermark_number(b["text"]):
                continue
            if is_footer(a["text"]) or is_footer(b["text"]):
                continue
            # 判斷 bounding box 是否重疊
            overlap_x = a["l"] < b["r"] - TOLERANCE and a["r"] > b["l"] + TOLERANCE
            overlap_y = a["t"] < b["b"] - TOLERANCE and a["b"] > b["t"] + TOLERANCE
            if overlap_x and overlap_y:
                issues.append({
                    "type": "overlap",
                    "page": page_num,
                    "name_a": a["name"],
                    "text_a": a["text"][:35],
                    "y_a": f"{a['t']:.2f}~{a['b']:.2f}",
                    "name_b": b["name"],
                    "text_b": b["text"][:35],
                    "y_b": f"{b['t']:.2f}~{b['b']:.2f}",
                })
    return issues

def validate(pptx_path: str) -> bool:
    prs = Presentation(pptx_path)
    W = prs.slide_width.inches
    H = prs.slide_height.inches
    all_issues = []

    for i, slide in enumerate(prs.slides):
        all_issues.extend(check_slide(slide, W, H, i + 1))

    if not all_issues:
        print(f"✓ 全部 {len(prs.slides)} 頁：無 overflow，無文字重疊")
        return True

    # 分類輸出
    overflows = [x for x in all_issues if x["type"].startswith("overflow")]
    overlaps  = [x for x in all_issues if x["type"] == "overlap"]

    if overflows:
        print(f"\n❌ Overflow（{len(overflows)} 個）：")
        for o in overflows:
            print(f"  第{o['page']}頁 [{o['name']}] \"{o['text']}\"")
            print(f"    → {o['detail']}")

    if overlaps:
        print(f"\n❌ 文字重疊（{len(overlaps)} 個）：")
        for o in overlaps:
            print(f"  第{o['page']}頁 [{o['name_a']}] \"{o['text_a']}\" (y={o['y_a']})")
            print(f"    ↕  [{o['name_b']}] \"{o['text_b']}\" (y={o['y_b']})")

    print(f"\n共 {len(all_issues)} 個問題需要修正。")
    return False

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python3 validate_pptx.py <path_to_pptx>")
        sys.exit(1)
    ok = validate(sys.argv[1])
    sys.exit(0 if ok else 1)
